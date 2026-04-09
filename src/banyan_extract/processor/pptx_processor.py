import re
import io
import PIL
import logging
import sys
import os
from typing import List, Tuple, Union

from tqdm import tqdm

# Use centralized logging
from ..utils.logging_config import get_logger

logger = get_logger(__name__)

try:
    from surya.texify import TexifyPredictor # May need to find an alternative
    USE_OCR = True
except ImportError as e:
    USE_OCR = False
    logger.warning(f"Surya OCR not installed; OCR not available for PPTX processing: {e}")

from pptx import Presentation
from pptx.shapes.group import GroupShape

from .processor import Processor
from ..output.pptx_output import PptxOutput
from ..ocr.nemotron_ocr import NemotronOCR


try:
    class MarkdownTexifyPredictor(TexifyPredictor):
        def fix_fences(self, text: str) -> str:
            text = re.sub(r'<math display="block">(.*?)</math>',r'$$\1$$', text, flags=re.DOTALL)
            text = re.sub(r'<math>(.*?)</math>',r'$\1$', text, flags=re.DOTALL)
            if re.search(r'<math display="block">', text):
                text = ""
            if re.search(r'<math>', text):
                text = ""
            return text
except ImportError as e:
    logger.warning(f"Failed to create MarkdownTexifyPredictor class: {e}")
    class MarkdownTexifyPredictor:
        pass
except Exception as e:
    logger.error(f"Unexpected error creating MarkdownTexifyPredictor: {e}")
    class MarkdownTexifyPredictor:
        pass


class PptxProcessor(Processor):
    
    def __init__(self, ocr_backend="surya", nemotron_endpoint="", nemotron_model="nvidia/nemoretriever-parse"):
        """
        Initialize PPTX processor with OCR backend selection.
        
        Args:
            ocr_backend: Which OCR backend to use ('surya' or 'nemotron')
            nemotron_endpoint: URL for Nemotron parse endpoint (if using nemotron)
            nemotron_model: Model name for Nemotron OCR
        """
        super().__init__()
        self.ocr_backend = None  # Single member variable for OCR backend
        self.ocr_available = False

        # Initialize selected OCR backend
        if ocr_backend == "nemotron":
            try:
                self.ocr_backend = NemotronOCR(
                    endpoint_url=nemotron_endpoint,
                    model_name=nemotron_model
                )
                self.ocr_available = True
            except ImportError as e:
                logger.warning(f"Nemotron OCR dependencies not available: {e}")
                logger.warning("To enable Nemotron OCR for PPTX processing, install nemotronparse dependencies: pip install .[nemotronparse]")
                self.ocr_available = False
            except ValueError as e:
                logger.error(f"Invalid Nemotron OCR configuration: {e}")
                self.ocr_available = False
            except Exception as e:
                logger.error(f"Failed to initialize Nemotron OCR: {e}")
                self.ocr_available = False
        else:  # Default to Surya
            try:
                self.ocr_backend = MarkdownTexifyPredictor()
                self.ocr_available = USE_OCR
            except ImportError as e:
                self.ocr_available = False
                logger.warning(f"Surya OCR not installed; OCR not available for PPTX processing: {e}")
                logger.warning("To enable Surya OCR for PPTX processing, install marker dependencies: pip install .[marker]")
            except Exception as e:
                self.ocr_available = False
                logger.error(f"Unexpected error initializing Surya OCR: {e}")
                logger.warning("Surya OCR initialization failed. PPTX processing will continue without OCR.")

    def ocr_image(self, image):
        """
        Perform OCR using the selected backend.
        
        Args:
            image: PIL Image object to perform OCR on
            
        Returns:
            Extracted text from the image, or empty string if OCR fails
            
        Raises:
            ValueError: If OCR backend is not available or image is invalid
        """
        if not self.ocr_available or self.ocr_backend is None:
            logger.debug("OCR not available")
            return ""

        if image is None:
            logger.warning("Cannot perform OCR on None image")
            return ""

        try:
            if isinstance(self.ocr_backend, NemotronOCR):
                # Use Nemotron OCR
                return self.ocr_backend.ocr_image(image)
            else:
                # Use Surya OCR
                ocr_output = self.ocr_backend([image])[0]
                if ocr_output.text is not None:
                    return ocr_output.text
                return ""
        except ValueError as e:
            logger.error(f"OCR validation error: {e}")
            return ""
        except AttributeError as e:
            logger.error(f"OCR backend attribute error: {e}")
            return ""
        except Exception as e:
            logger.error(f"OCR failed: {e}")
            return ""

    def process_image(self, image):
        if 'wmf' in image.content_type:
            print("pptx contains a wmf image, skipping image")
            return None
        else:
            byte_io = io.BytesIO(image.blob)
            image_object = PIL.Image.open(byte_io)
            return image_object

    def process_document(self, filepath, rotation_angle: Union[int, float] = 0,
                       auto_detect_rotation: bool = False, 
                       rotation_confidence_threshold: float = 0.7):
        """
        Process a single PPTX document.
        
        Args:
            filepath: Path to the PPTX document file
            rotation_angle: Rotation angle in degrees (default: 0)
            auto_detect_rotation: Whether to automatically detect rotation (default: False)
            rotation_confidence_threshold: Minimum confidence for auto rotation detection (default: 0.7)
            
        Returns:
            PptxOutput object containing processed document data
            
        Raises:
            FileNotFoundError: If the input file cannot be found
            PermissionError: If the input file cannot be read
            ValueError: If the file is not a valid PPTX or processing fails
            Exception: For other processing errors
            
        Note:
            Rotation is not currently supported for PptxProcessor as it works with 
            structured slides. For future implementation, we would need to rotate 
            individual images within slides.
        """
        # Note: PPTX processor doesn't currently support rotation as it works with structured slides
        # For future implementation, we would need to rotate individual images within slides
        if rotation_angle != 0 or auto_detect_rotation:
            logger.warning(f"Rotation is not currently supported for PptxProcessor. "
                          f"Angle {rotation_angle} and auto-detection will be ignored.")
        
        try:
            prs = Presentation(filepath)
            images = []
            slide_texts = []
            for slide in prs.slides:
                slide_text = []
                images.append([])
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for par in shape.text_frame.paragraphs:
                            slide_text.append(par.text)
                    if "PIC" in str(shape.shape_type):
                        try:
                            image = self.process_image(shape.image)
                            if image is not None:
                                image_ocr = self.ocr_image(image)
                                slide_text.append(image_ocr)
                                images[-1].append(image)
                        except Exception as img_e:
                            logger.warning(f"Failed to process image in slide: {img_e}")
                            continue
                    if "GROUP" in str(shape.shape_type):
                        for sub_shape in shape.shapes:
                            if "PIC" in str(sub_shape.shape_type):
                                try:
                                    image = self.process_image(sub_shape.image)
                                    if image is not None:
                                        image_ocr = self.ocr_image(image)
                                        slide_text.append(image_ocr)
                                        images[-1].append(image)
                                except Exception as img_e:
                                    logger.warning(f"Failed to process grouped image in slide: {img_e}")
                                    continue

                            if sub_shape.has_text_frame:
                                for par in sub_shape.text_frame.paragraphs:
                                    slide_text.append(par.text)
                slide_texts.append("\n".join(slide_text))

            metadata = {}
            return PptxOutput(slide_texts, images, metadata)
        except FileNotFoundError as e:
            logger.error(f"PPTX file not found: {filepath}")
            raise FileNotFoundError(f"PPTX file not found: {filepath}") from e
        except PermissionError as e:
            logger.error(f"Permission denied reading PPTX file: {filepath}")
            raise PermissionError(f"Permission denied reading PPTX file: {filepath}") from e
        except ValueError as e:
            logger.error(f"Invalid PPTX file format: {filepath}")
            raise ValueError(f"Invalid PPTX file format: {filepath}") from e
        except Exception as e:
            logger.error(f"Error processing PPTX document {filepath}: {e}")
            raise ValueError(f"Failed to process PPTX document {filepath}: {e}") from e

    def process_batch_documents(self, filepaths, rotation_angle: Union[int, float] = 0,
                               auto_detect_rotation: bool = False, 
                               rotation_confidence_threshold: float = 0.7):
        """
        Process multiple PPTX documents in batch.
        
        Args:
            filepaths: List of paths to PPTX document files
            rotation_angle: Rotation angle in degrees (default: 0)
            auto_detect_rotation: Whether to automatically detect rotation (default: False)
            rotation_confidence_threshold: Minimum confidence for auto rotation detection (default: 0.7)
            
        Returns:
            List of PptxOutput objects for each processed file
            
        Raises:
            FileNotFoundError: If any input file cannot be found
            PermissionError: If any input file cannot be read
            ValueError: If any file is not a valid PPTX or processing fails
            Exception: For other processing errors
            
        Note:
            Batch processing is not currently fully implemented for PPTX processor.
            Rotation is not currently supported for PptxProcessor.
        """
        # Note: Batch processing not currently implemented for PPTX processor
        if rotation_angle != 0 or auto_detect_rotation:
            logger.warning(f"Rotation is not currently supported for PptxProcessor. "
                          f"Angle {rotation_angle} and auto-detection will be ignored.")
        
         # For future implementation, we would process each file with rotation support
        file_outputs = []
        for filepath in filepaths:
            try:
                output = self.process_document(filepath)
                file_outputs.append(output)
            except FileNotFoundError as e:
                logger.error(f"PPTX file not found in batch: {filepath}")
                raise FileNotFoundError(f"PPTX file not found in batch: {filepath}") from e
            except PermissionError as e:
                logger.error(f"Permission denied reading PPTX file in batch: {filepath}")
                raise PermissionError(f"Permission denied reading PPTX file in batch: {filepath}") from e
            except ValueError as e:
                logger.error(f"Invalid PPTX file format in batch: {filepath}")
                raise ValueError(f"Invalid PPTX file format in batch: {filepath}") from e
            except Exception as e:
                logger.error(f"Failed to process PPTX file {filepath} in batch: {e}")
                raise ValueError(f"Failed to process PPTX file {filepath} in batch: {e}") from e
        
        return file_outputs
