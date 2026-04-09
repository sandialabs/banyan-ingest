# Integration tests for PPTX processor backends
# Tests both Nemotron and Surya backends with real PPTX files

import pytest
from banyan_extract.processor.pptx_processor import PptxProcessor


class TestPPTXBackendsIntegration:
    """Integration tests for PPTX processor with different backends using real files."""

    def test_nemotron_backend_integration(self, sample_pptx_file):
        """Test Nemotron backend integration with actual PPTX processing."""
        # Test with Nemotron backend (will use real OCR if available, or skip gracefully)
        processor = PptxProcessor(ocr_backend="nemotron")
        
        # Process the actual PPTX file
        output = processor.process_document(str(sample_pptx_file))
        
        # Verify output structure
        assert output is not None
        assert hasattr(output, 'text')
        assert hasattr(output, 'images')
        assert hasattr(output, 'metadata')
        
        # Verify text was extracted from text frames
        assert len(output.text) > 0
        
        # Verify we have text content from the slides
        # The sample file should have content like "This is a title", "This is a subtitle", etc.
        all_text = " ".join(output.text)
        assert len(all_text.strip()) > 0, "Should extract text from PPTX slides"

    def test_surya_backend_integration(self, sample_pptx_file):
        """Test Surya backend integration with actual PPTX processing."""
        # Skip test if surya dependencies are not available
        try:
            from surya.texify import TexifyPredictor
        except ImportError:
            pytest.skip("Surya dependencies not available, skipping Surya backend test")
        
        # Test with Surya backend
        processor = PptxProcessor(ocr_backend="surya")
        
        # Process the actual PPTX file
        output = processor.process_document(str(sample_pptx_file))
        
        # Verify output structure
        assert output is not None
        assert hasattr(output, 'text')
        assert hasattr(output, 'images')
        assert hasattr(output, 'metadata')
        
        # Verify text was extracted from text frames
        assert len(output.text) > 0
        
        # Verify we have text content from the slides
        all_text = " ".join(output.text)
        assert len(all_text.strip()) > 0, "Should extract text from PPTX slides"

    def test_no_backend_integration(self, sample_pptx_file):
        """Test processing with no OCR backend available."""
        # Test with unknown backend (no OCR)
        processor = PptxProcessor(ocr_backend="unknown")
        
        # Process the actual PPTX file
        output = processor.process_document(str(sample_pptx_file))
        
        # Verify output structure
        assert output is not None
        assert hasattr(output, 'text')
        assert hasattr(output, 'images')
        assert hasattr(output, 'metadata')
        
        # Verify text was extracted (text frames should still work)
        assert len(output.text) > 0
        
        # Verify we have text content from the slides
        all_text = " ".join(output.text)
        assert len(all_text.strip()) > 0, "Should extract text from PPTX slides even without OCR"

    def test_backend_switching_integration(self, sample_pptx_file):
        """Test switching between backends during processing."""
        # First process with Nemotron
        nemotron_processor = PptxProcessor(ocr_backend="nemotron")
        nemotron_output = nemotron_processor.process_document(str(sample_pptx_file))
        
        # Verify Nemotron processing
        assert nemotron_output is not None
        assert len(nemotron_output.text) > 0
        
        # Then process with Surya (if available)
        try:
            from surya.texify import TexifyPredictor
            surya_processor = PptxProcessor(ocr_backend="surya")
            surya_output = surya_processor.process_document(str(sample_pptx_file))
            
            # Verify Surya processing
            assert surya_output is not None
            assert len(surya_output.text) > 0
            
            # Verify both backends produced valid output with same number of slides
            assert len(nemotron_output.text) == len(surya_output.text)
        except ImportError:
            # Surya not available, just test Nemotron worked
            assert nemotron_output is not None
            assert len(nemotron_output.text) > 0

    def test_backend_consistency_across_runs(self, sample_pptx_file):
        """Test that backend produces consistent results across multiple runs."""
        # Test consistency with Nemotron backend
        processor = PptxProcessor(ocr_backend="nemotron")
        
        # Process the same file multiple times
        output1 = processor.process_document(str(sample_pptx_file))
        output2 = processor.process_document(str(sample_pptx_file))
        output3 = processor.process_document(str(sample_pptx_file))
        
        # Verify consistent results (same number of slides and structure)
        assert len(output1.text) == len(output2.text) == len(output3.text)
        assert len(output1.images) == len(output2.images) == len(output3.images)
        
        # Text content should be identical for text frames
        assert output1.text == output2.text == output3.text

    def test_backend_with_custom_configuration(self, sample_pptx_file):
        """Test backend with custom configuration."""
        # Test with custom configuration (will work even if Nemotron is not available)
        custom_endpoint = "http://custom-endpoint:8000"
        custom_model = "custom-model-name"
        
        processor = PptxProcessor(
            ocr_backend="nemotron",
            nemotron_endpoint=custom_endpoint,
            nemotron_model=custom_model
        )
        
        # Process the file
        output = processor.process_document(str(sample_pptx_file))
        
        # Verify processing works with custom configuration
        assert output is not None
        assert len(output.text) > 0


class TestPPTXBackendsEdgeCases:
    """Edge case tests for PPTX backends integration."""

    def test_backend_with_empty_pptx(self, tmp_path):
        """Test backend processing of empty PPTX file."""
        # Create an empty PPTX file
        from pptx import Presentation
        empty_pptx = tmp_path / "empty.pptx"
        
        prs = Presentation()
        # Add an empty slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        prs.save(str(empty_pptx))
        
        # Test with Nemotron backend
        processor = PptxProcessor(ocr_backend="nemotron")
        output = processor.process_document(str(empty_pptx))
        
        # Verify empty PPTX is handled gracefully
        assert output is not None
        assert len(output.text) >= 1  # At least one empty slide

    def test_backend_with_text_only_pptx(self, tmp_path):
        """Test backend processing of text-only PPTX file."""
        from pptx import Presentation
        text_only_pptx = tmp_path / "text_only.pptx"
        
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        
        # Add text content
        title = slide.shapes.title
        title.text = "Test Title"
        
        content = slide.placeholders[1]
        content.text = "Test Content"
        
        prs.save(str(text_only_pptx))
        
        # Test with Nemotron backend
        processor = PptxProcessor(ocr_backend="nemotron")
        output = processor.process_document(str(text_only_pptx))
        
        # Verify text-only PPTX is handled correctly
        assert output is not None
        assert len(output.text) >= 1
        assert "Test Title" in output.text[0] or "Test Content" in output.text[0]

    def test_backend_with_image_only_pptx(self, tmp_path):
        """Test backend processing of image-only PPTX file."""
        from pptx import Presentation
        from pptx.util import Inches
        from PIL import Image
        
        image_only_pptx = tmp_path / "image_only.pptx"
        
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        
        # Add an image
        # Create a simple test image
        test_image = Image.new('RGB', (100, 100), color='white')
        test_image_path = tmp_path / "test_image.png"
        test_image.save(str(test_image_path))
        
        slide.shapes.add_picture(str(test_image_path), Inches(1), Inches(1))
        
        prs.save(str(image_only_pptx))
        
        # Test with Nemotron backend
        processor = PptxProcessor(ocr_backend="nemotron")
        output = processor.process_document(str(image_only_pptx))
        
        # Verify image-only PPTX is handled correctly
        assert output is not None
        assert len(output.text) >= 1
        
        # If OCR is available, the image should be processed
        # If not, we should still have an empty slide
        assert len(output.text) >= 1  # At least some text should be present

    def test_real_pptx_file_content(self, sample_pptx_file):
        """Test that we can extract real content from the sample PPTX file."""
        # Test with Nemotron backend to extract real content
        processor = PptxProcessor(ocr_backend="nemotron")
        output = processor.process_document(str(sample_pptx_file))
        
        # Verify we extracted real content from the sample file
        assert output is not None
        assert len(output.text) > 0
        
        # The sample file should contain specific expected content
        all_text = " ".join(output.text)
        
        # Check for expected content from the slides.pptx file
        expected_phrases = ["This is a title", "This is a subtitle", "This is another slide"]
        found_phrases = [phrase for phrase in expected_phrases if phrase in all_text]
        
        # We should find at least some of the expected content
        assert len(found_phrases) > 0, f"Expected to find some of {expected_phrases} in extracted text, but found none"
        
        # Verify we have the correct number of slides
        assert len(output.text) == 2, f"Expected 2 slides, got {len(output.text)}"


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
