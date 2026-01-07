import re
import io
import PIL
from typing import List, Tuple

from tqdm import tqdm

try:
    from surya.texify import TexifyPredictor # May need to find an alternative
    USE_OCR = True
except:
    USE_OCR = False
    print("Surya ocr not installed; ocr not available for pptx processing")

from pptx import Presentation
from pptx.shapes.group import GroupShape

from .processor import Processor
from ..output.pptx_output import PptxOutput


try:
    class MarkdownTexifyPredictor(TexifyPredictor):
        def fix_fences(self, text: str) -> str:
            text = re.sub(r'<math display="block">(.*?)</math>',r'$$\1$$', text, flags=re.DOTALL)
            text = re.sub(r'<math>(.*?)</math>',r'$\1$', text, flags=re.DOTALL)
            if re.search(r'<math display="block">', text):
                text = ""
            if re.search(r'<math>', text):
                text = ""
            return text
except:
    class MarkdownTexifyPredictor:
        pass


class PptxProcessor(Processor):
    
    def __init__(self):
        super().__init__()
        if USE_OCR:
            self.latex_predictor = MarkdownTexifyPredictor()
        else:
            self.latex_predictor = None

    def ocr_image(self, image):
        if USE_OCR:
            ocr_output = self.latex_predictor([image])[0]
            if ocr_output.text is not None:
                return ocr_output.text
        return ""

    def process_image(self, image):
        if 'wmf' in image.content_type:
            print("pptx contains a wmf image, skipping image")
            return None
        else:
            byte_io = io.BytesIO(image.blob)
            image_object = PIL.Image.open(byte_io)
            return image_object

    def process_document(self, filepath):
        prs = Presentation(filepath)
        images = []
        slide_texts = []
        for slide in prs.slides:
            slide_text = []
            images.append([])
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for par in shape.text_frame.paragraphs:
                        slide_text.append(par.text)
                if "PIC" in str(shape.shape_type):
                    image = self.process_image(shape.image)
                    if image is not None:
                        image_ocr = self.ocr_image(image)
                        slide_text.append(image_ocr)
                        images[-1].append(image)
                if "GROUP" in str(shape.shape_type):
                    for sub_shape in shape.shapes:
                        if "PIC" in str(sub_shape.shape_type):
                            image = self.process_image(sub_shape.image)
                            if image is not None:
                                image_ocr = self.ocr_image(image)
                                slide_text.append(image_ocr)
                                images[-1].append(image)

                        if sub_shape.has_text_frame:
                            for par in sub_shape.text_frame.paragraphs:
                                slide_text.append(par.text)
            slide_texts.append("\n".join(slide_text))

        metadata = {}
        return PptxOutput(slide_texts, images, metadata)

    def process_batch_documents(self, filepaths):
        pass
